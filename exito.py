import streamlit as st
import requests
from urllib.parse import urlparse
import matplotlib.pyplot as plt
import re
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

# ============================================
# Funciones Auxiliares
# ============================================

def parse_number(text):
    """
    Función para parsear números con diferentes formatos.
    Convierte cadenas de texto a números flotantes.
    """
    text = re.sub(r'[^\d.,]', '', text)
    if text.count(',') > text.count('.'):
        text = text.replace('.', '').replace(',', '.')
    else:
        text = text.replace(',', '')
    try:
        return float(text)
    except ValueError:
        return None

@st.cache_data(show_spinner=False)
def obtener_busqueda_serper(query, api_key):
    """
    Función para obtener resultados de búsqueda utilizando la API de Serper.
    """
    serper_url = "https://google.serper.dev/search"
    headers_serper = {
        "X-API-KEY": api_key,
        "Content-Type": "application/json"
    }
    data_serper = {"q": query}
    response_serper = requests.post(serper_url, headers=headers_serper, json=data_serper, timeout=10)
    response_serper.raise_for_status()
    search_results = response_serper.json()
    snippets = []
    if "organic" in search_results:
        for item in search_results["organic"]:
            if "snippet" in item:
                snippets.append(item["snippet"])
    return "\n\n".join(snippets) if snippets else "No se encontraron resultados relevantes."

@st.cache_data(show_spinner=False)
def obtener_analisis_together(search_summary, api_key):
    """
    Función para obtener análisis detallado utilizando la API de Together.
    """
    together_url = "https://api.together.xyz/v1/chat/completions"
    headers_together = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    messages = [
        {
            "role": "system",
            "content": (
                "Eres un experto en análisis de plataformas digitales con un enfoque en las demandas del mercado actual. "
                "Proporciona una evaluación detallada del potencial de éxito de la plataforma digital basada en la siguiente información. "
                "Incluye recomendaciones sobre aspectos de forma (diseño, usabilidad, interfaz) y fondo (funcionalidades, contenido, estrategia de mercado), señalando lo que sobra y lo que falta. "
                "Además, expresa el potencial de éxito en términos de porcentaje, proporciona una estimación del máximo de visitantes al día y un resumen ejecutivo de los hallazgos clave."
                "\n\n"
                "Nota: La estimación de visitantes diarios se basa en la versión mejorada de la plataforma, incorporando los cambios sugeridos."
            )
        },
        {
            "role": "user",
            "content": f"Aquí hay información sobre la plataforma: {search_summary}"
        }
    ]
    data_together = {
        "model": "mistralai/Mixtral-8x7B-Instruct-v0.1",
        "messages": messages,
        "max_tokens": 2512,
        "temperature": 0.7,
        "top_p": 0.7,
        "top_k": 50,
        "repetition_penalty": 1,
        "stop": ["<|eot_id|>"],
        "stream": False
    }

    response_together = requests.post(together_url, headers=headers_together, json=data_together, timeout=30)
    response_together.raise_for_status()
    analysis = response_together.json()
    if "choices" in analysis and len(analysis["choices"]) > 0:
        return analysis["choices"][0]["message"]["content"]
    else:
        raise ValueError("Respuesta inesperada de Together API.")

def generar_graficas(secciones):
    """
    Función para generar las gráficas necesarias basadas en las secciones del análisis.
    Retorna un diccionario con el título de la sección y el buffer de la imagen.
    """
    plots = {}
    # Potencial de Éxito
    if "Potencial de Éxito" in secciones:
        porcentaje_val = parse_number(secciones["Potencial de Éxito"])
        if porcentaje_val is not None:
            fig, ax = plt.subplots(figsize=(2, 2))
            ax.pie([porcentaje_val, 100 - porcentaje_val], colors=['#4CAF50', '#CCCCCC'], startangle=90, counterclock=False)
            ax.axis('equal')
            buf = BytesIO()
            plt.savefig(buf, format='png', bbox_inches='tight')
            plt.close(fig)
            buf.seek(0)
            plots["Potencial de Éxito"] = buf

    # Estimación de Visitantes Diarios
    if "Estimación de Visitantes Diarios" in secciones:
        est_visitors_val = parse_number(secciones["Estimación de Visitantes Diarios"])
        if est_visitors_val is not None:
            est_visitors_val = int(est_visitors_val)
            fig, ax = plt.subplots(figsize=(4, 1))
            ax.barh([''], [est_visitors_val], color='#4CAF50')
            ax.set_xlim(0, est_visitors_val * 1.2)
            ax.set_xlabel('Número de Visitantes')
            ax.set_yticks([])
            plt.tight_layout()
            buf = BytesIO()
            plt.savefig(buf, format='png', bbox_inches='tight')
            plt.close(fig)
            buf.seek(0)
            plots["Estimación de Visitantes Diarios"] = buf

    return plots

def generar_pptx(secciones, plots):
    """
    Función para generar una presentación de PowerPoint (PPTX) con el análisis y las gráficas.
    Divide el análisis en varias diapositivas, una por cada sección.
    Retorna un buffer con el archivo PPTX.
    """
    prs = Presentation()
    
    # Diseño de diapositiva: Título y contenido
    slide_layout_title = prs.slide_layouts[0]  # Título
    slide_layout_content = prs.slide_layouts[1]  # Título y contenido

    # Diapositiva de Título
    slide_title = prs.slides.add_slide(slide_layout_title)
    title = slide_title.shapes.title
    title.text = "Análisis de Potencial de Éxito"

    # Iterar sobre las secciones y añadir diapositivas
    for titulo, contenido in secciones.items():
        slide = prs.slides.add_slide(slide_layout_content)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = titulo

        # Agregar el contenido de la sección
        tf = body.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = contenido
        p.font.size = Pt(12)

        # Añadir gráfica si existe
        if titulo in plots:
            img_buffer = plots[titulo]
            img = Image.open(img_buffer)
            img_path = f"{titulo}.png"
            img.save(img_path)

            # Insertar la imagen en la diapositiva
            left = Inches(1)
            top = Inches(3)
            height = Inches(3)
            slide.shapes.add_picture(img_path, left, top, height=height)

    # Guardar la presentación en un buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# ============================================
# Aplicación Streamlit
# ============================================

def main():
    # Título de la aplicación
    st.title("📈 Análisis de Potencial de Éxito de Plataformas Digitales")

    # Descripción de la aplicación
    st.markdown("""
    Esta aplicación analiza el potencial de éxito de una plataforma digital basada en su URL. Utiliza las APIs de Serper para obtener información relevante sobre la plataforma y de Together para evaluar su potencial en el mercado actual. Además, proporciona recomendaciones detalladas para mejorar tanto en forma como en contenido, así como una estimación del máximo de visitantes diarios que puede recibir la plataforma.
    """)

    # Entrada de la URL
    url_input = st.text_input("🔗 Ingresa la URL de la plataforma digital que deseas analizar:", "")

    # Botón para iniciar el análisis
    if st.button("✅ Analizar"):
        if not url_input:
            st.error("⚠️ Por favor, ingresa una URL válida.")
        else:
            # Validar y formatear la URL
            parsed_url = urlparse(url_input)
            if not parsed_url.scheme:
                url_input = "https://" + url_input  # Preferir HTTPS
                parsed_url = urlparse(url_input)

            domain = parsed_url.netloc
            if not domain:
                st.error("⚠️ URL inválida. Por favor, intenta nuevamente.")
                st.stop()

            try:
                response = requests.get(url_input, timeout=10)
                response.raise_for_status()
            except requests.exceptions.HTTPError as http_err:
                st.error(f"⚠️ Error HTTP al acceder a la URL: {http_err}")
                st.stop()
            except requests.exceptions.ConnectionError:
                st.error("⚠️ Error de conexión. Verifica tu red y la URL ingresada.")
                st.stop()
            except requests.exceptions.Timeout:
                st.error("⚠️ Tiempo de espera excedido al intentar acceder a la URL.")
                st.stop()
            except requests.exceptions.RequestException as e:
                st.error(f"⚠️ Error al acceder a la URL: {e}")
                st.stop()

            st.info("🔄 Procesando la URL...")

            # Realizar búsqueda con Serper API
            serper_api_key = st.secrets["serper_api_key"]
            query = f"Información sobre {domain}"

            with st.spinner("🔍 Realizando búsqueda con Serper..."):
                try:
                    search_summary = obtener_busqueda_serper(query, serper_api_key)
                except requests.exceptions.HTTPError as http_err:
                    st.error(f"❌ Error HTTP al acceder a Serper API: {http_err}")
                    st.stop()
                except requests.exceptions.RequestException as e:
                    st.error(f"❌ Error al acceder a Serper API: {e}")
                    st.stop()

            # Preparar el mensaje para Together API
            together_api_key = st.secrets["together_api_key"]

            with st.spinner("🧠 Analizando con Together..."):
                try:
                    result = obtener_analisis_together(search_summary, together_api_key)
                except requests.exceptions.HTTPError as http_err:
                    st.error(f"❌ Error HTTP al acceder a Together API: {http_err}")
                    st.stop()
                except requests.exceptions.RequestException as e:
                    st.error(f"❌ Error al acceder a Together API: {e}")
                    st.stop()
                except ValueError as ve:
                    st.error(f"❌ {ve}")
                    st.stop()

            # Procesar y unificar el análisis
            # Separar el análisis en secciones utilizando títulos en negrita
            secciones = {}
            current_section = None
            for line in result.split('\n'):
                line = line.strip()
                match = re.match(r'\*\*(.*?)\*\*:', line)
                if match:
                    section_title = match.group(1).strip()
                    secciones[section_title] = ""
                    current_section = section_title
                elif current_section:
                    secciones[current_section] += line + "\n"

            # Manejar caso sin secciones
            if not secciones:
                secciones["Contenido"] = result

            # Generar las gráficas y obtener los buffers de imagen
            plots = generar_graficas(secciones)

            # Crear un contenedor para unificar el análisis
            with st.container():
                st.subheader("📊 Análisis Unificado")
                for titulo, contenido in secciones.items():
                    st.markdown(f"### 📌 **{titulo}**")
                    st.write(contenido)

                    # Incluir visualizaciones dentro de las secciones pertinentes
                    if titulo in plots:
                        # Mostrar la gráfica en Streamlit
                        st.image(plots[titulo], use_column_width=True)

                        # Agregar explicación detallada para la estimación de visitantes diarios
                        if titulo == "Estimación de Visitantes Diarios":
                            st.markdown("""
                            **¿Cómo se estima el número máximo de visitantes diarios?**
                            
                            La estimación de visitantes diarios se basa en un análisis detallado de la plataforma digital, considerando varios factores clave que influyen en su capacidad de atraer y retener usuarios. A continuación, se detallan los aspectos considerados para esta estimación:
                            
                            1. **Optimización de la Usabilidad y Diseño:**
                               - **Mejoras en la Interfaz de Usuario (UI):** Un diseño intuitivo y atractivo facilita la navegación y mejora la experiencia del usuario, lo que puede incrementar la retención y la recomendación boca a boca.
                               - **Responsive Design:** La adaptación del sitio web a diferentes dispositivos (móviles, tabletas, computadoras) asegura que una mayor audiencia pueda acceder y utilizar la plataforma sin inconvenientes.
                            
                            2. **Contenido Relevante y de Calidad:**
                               - **Actualización Regular del Contenido:** Mantener el contenido fresco y actualizado atrae a usuarios recurrentes y mejora el posicionamiento en motores de búsqueda.
                               - **SEO (Search Engine Optimization):** La optimización para motores de búsqueda incrementa la visibilidad de la plataforma, atrayendo más tráfico orgánico.
                            
                            3. **Estrategias de Marketing y Promoción:**
                               - **Campañas de Marketing Digital:** Utilizar canales como redes sociales, correo electrónico y publicidad pagada puede aumentar significativamente el tráfico hacia la plataforma.
                               - **Colaboraciones y Alianzas Estratégicas:** Asociarse con otras empresas o influencers puede ampliar la base de usuarios potenciales.
                            
                            4. **Funcionalidades Mejoradas:**
                               - **Integración de Funcionalidades Clave:** Incorporar herramientas y características que satisfagan las necesidades de los usuarios puede aumentar el tiempo de permanencia y la frecuencia de visitas.
                               - **Análisis de Datos y Personalización:** Utilizar datos para personalizar la experiencia del usuario mejora la satisfacción y fomenta la lealtad.
                            
                            5. **Soporte y Atención al Cliente:**
                               - **Canales de Soporte Eficientes:** Ofrecer soporte rápido y efectivo resuelve problemas de usuarios y mejora la percepción general de la plataforma.
                            
                            **Conclusión:**
                            
                            Al implementar las mejoras sugeridas en estos aspectos, se espera que la plataforma digital no solo aumente su atractivo y funcionalidad, sino que también mejore su capacidad para atraer y retener un mayor número de visitantes diarios. La combinación de un diseño optimizado, contenido de calidad, estrategias de marketing efectivas y funcionalidades avanzadas contribuye significativamente a la estimación del número máximo de visitantes diarios.
                            """)

            # Generar la presentación PPTX
            with st.spinner("📊 Generando presentación PPTX..."):
                pptx_buffer = generar_pptx(secciones, plots)

            # Botón para descargar el PPTX
            st.download_button(
                label="📥 Descargar Análisis en PPTX",
                data=pptx_buffer,
                file_name="analisis_plataforma.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

# Ejecutar la aplicación
if __name__ == "__main__":
    main()
